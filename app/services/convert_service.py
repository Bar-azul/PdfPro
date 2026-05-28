"""
ConvertService
==============
Handles all file format conversions:
  PDF → Word, Excel, PowerPoint, JPG/PNG
  Word / Excel / PowerPoint / Images → PDF
"""

import logging
import os
import shutil
import subprocess
import tempfile
import time
import uuid
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image

from ..config import settings
from ..services.pdf_service import PDFService, _temp_pdf, _temp_file, _ms

logger = logging.getLogger(__name__)


OFFICE_EXTENSIONS = {
    ".doc",
    ".docx",
    ".odt",
    ".rtf",
    ".xls",
    ".xlsx",
    ".ods",
    ".ppt",
    ".pptx",
    ".odp",
}


def _resolve_libreoffice_path() -> str:
    """
    Resolve LibreOffice / soffice executable path.

    Priority:
    1. settings.LIBREOFFICE_PATH
    2. Common Windows installation path
    3. soffice from PATH
    4. libreoffice from PATH
    """

    configured_path = getattr(settings, "LIBREOFFICE_PATH", None)

    if configured_path:
        configured_path = str(configured_path).strip()

        if configured_path and Path(configured_path).exists():
            return configured_path

        # If configured as command name, allow it.
        found_configured = shutil.which(configured_path)
        if found_configured:
            return found_configured

    common_windows_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
    if Path(common_windows_path).exists():
        return common_windows_path

    soffice_path = shutil.which("soffice")
    if soffice_path:
        return soffice_path

    libreoffice_path = shutil.which("libreoffice")
    if libreoffice_path:
        return libreoffice_path

    raise RuntimeError(
        "LibreOffice executable was not found. "
        "Install LibreOffice or set LIBREOFFICE_PATH to the full path of soffice.exe. "
        r'Example for Windows: C:\Program Files\LibreOffice\program\soffice.exe'
    )


def _safe_unlink(path: Path) -> None:
    try:
        path.unlink(missing_ok=True)
    except Exception:
        logger.warning("Failed to delete file: %s", path, exc_info=True)


class ConvertService:

    # ────────────────────────────────────────────────────────────────────────
    # PDF → Other formats
    # ────────────────────────────────────────────────────────────────────────

    @staticmethod
    def pdf_to_word(pdf_path: Path) -> Path:
        """Convert PDF to .docx using pdf2docx."""
        from pdf2docx import Converter

        t0 = time.time()
        pdf_path = Path(pdf_path).resolve()

        if not pdf_path.exists():
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")

        out = _temp_file("converted", ".docx")

        cv = Converter(str(pdf_path))
        try:
            cv.convert(str(out), start=0, end=None)
        finally:
            cv.close()

        logger.info("PDF→Word in %sms → %s", _ms(t0), out)
        return out

    @staticmethod
    def pdf_to_excel(pdf_path: Path) -> Path:
        """
        Convert PDF tables to .xlsx using pdfplumber for table extraction.
        Falls back to text extraction if no tables found.
        """
        import pdfplumber
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side

        t0 = time.time()
        pdf_path = Path(pdf_path).resolve()

        if not pdf_path.exists():
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")

        out = _temp_file("converted", ".xlsx")

        wb = openpyxl.Workbook()
        wb.remove(wb.active)

        header_font = Font(bold=True, color="FFFFFF", name="Calibri", size=11)
        header_fill = PatternFill(fill_type="solid", fgColor="0D1B2A")
        border = Border(
            left=Side(style="thin"),
            right=Side(style="thin"),
            top=Side(style="thin"),
            bottom=Side(style="thin"),
        )

        found_any_table = False

        with pdfplumber.open(str(pdf_path)) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                tables = page.extract_tables()

                if tables:
                    found_any_table = True

                    for tbl_num, table in enumerate(tables, start=1):
                        ws_name = f"עמ' {page_num}"
                        if tbl_num > 1:
                            ws_name += f" טבלה {tbl_num}"

                        ws = wb.create_sheet(title=ws_name[:31])

                        for row_idx, row in enumerate(table, start=1):
                            for col_idx, cell in enumerate(row, start=1):
                                c = ws.cell(
                                    row=row_idx,
                                    column=col_idx,
                                    value=cell or "",
                                )
                                c.border = border
                                c.alignment = Alignment(
                                    wrap_text=True,
                                    horizontal="right",
                                )

                                if row_idx == 1:
                                    c.font = header_font
                                    c.fill = header_fill

                        for col in ws.columns:
                            max_len = max(
                                (
                                    len(str(c.value))
                                    if c.value is not None
                                    else 0
                                    for c in col
                                ),
                                default=10,
                            )
                            ws.column_dimensions[col[0].column_letter].width = min(
                                max_len + 4,
                                60,
                            )

                else:
                    text = page.extract_text() or ""
                    ws = wb.create_sheet(title=f"עמ' {page_num} טקסט"[:31])

                    if text.strip():
                        for line_num, line in enumerate(text.splitlines(), start=1):
                            ws.cell(row=line_num, column=1, value=line)
                    else:
                        ws.cell(row=1, column=1, value="לא נמצא טקסט בעמוד זה")

            if not wb.worksheets:
                ws = wb.create_sheet(title="ריק")
                ws.cell(row=1, column=1, value="לא נמצאו טבלאות במסמך")

        wb.save(str(out))

        logger.info(
            "PDF→Excel in %sms, found tables: %s → %s",
            _ms(t0),
            found_any_table,
            out,
        )
        return out

    @staticmethod
    def pdf_to_pptx(pdf_path: Path, dpi: int = 150) -> Path:
        """
        Convert PDF to .pptx:
        render each page as an image, then insert into slides.
        Preserves visual appearance.
        """
        from pptx import Presentation

        t0 = time.time()
        pdf_path = Path(pdf_path).resolve()

        if not pdf_path.exists():
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")

        image_paths = PDFService.render_pages(pdf_path, dpi=dpi, fmt="png")

        if not image_paths:
            raise RuntimeError("PDF rendering produced no images")

        prs = Presentation()

        with fitz.open(pdf_path) as doc:
            if doc.page_count == 0:
                raise RuntimeError("PDF has no pages")

            first_page = doc[0].rect
            width_emu = int(first_page.width / 72 * 914400)
            height_emu = int(first_page.height / 72 * 914400)
            prs.slide_width = width_emu
            prs.slide_height = height_emu

        blank_layout = prs.slide_layouts[6]

        try:
            for img_path in image_paths:
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    str(img_path),
                    left=0,
                    top=0,
                    width=prs.slide_width,
                    height=prs.slide_height,
                )
        finally:
            for img_path in image_paths:
                _safe_unlink(Path(img_path))

        out = _temp_file("converted", ".pptx")
        prs.save(str(out))

        logger.info("PDF→PPTX in %sms, slides: %s → %s", _ms(t0), len(image_paths), out)
        return out

    @staticmethod
    def pdf_to_images(
        pdf_path: Path,
        dpi: int = 150,
        fmt: str = "jpg",
        quality: int = 85,
        pages: list[int] | None = None,
    ) -> list[Path]:
        """Convert PDF pages to image files. Returns list of paths."""
        pdf_path = Path(pdf_path).resolve()

        if not pdf_path.exists():
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")

        return PDFService.render_pages(
            pdf_path,
            dpi=dpi,
            fmt=fmt,
            quality=quality,
            pages=pages,
        )

    @staticmethod
    def pdf_to_text(pdf_path: Path) -> str:
        """Extract all text from PDF."""
        pdf_path = Path(pdf_path).resolve()

        if not pdf_path.exists():
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")

        with fitz.open(pdf_path) as doc:
            parts = [page.get_text() for page in doc]

        return "\n\n--- עמוד חדש ---\n\n".join(parts)

    # ────────────────────────────────────────────────────────────────────────
    # Other formats → PDF
    # ────────────────────────────────────────────────────────────────────────

    @staticmethod
    def office_to_pdf(input_path: Path) -> Path:
        """
        Convert Word / Excel / PowerPoint to PDF using LibreOffice.

        Important:
        On Windows, LibreOffice sometimes fails with:
        "Error: source file could not be loaded"

        To avoid that:
        1. Resolve the input path to an absolute path.
        2. Validate that the file exists and is not empty.
        3. Copy the source file into a temporary folder.
        4. Rename it to a clean ASCII-only filename.
        5. Run LibreOffice on the safe copied file.
        """
        t0 = time.time()
        input_path = Path(input_path).resolve()

        if not input_path.exists():
            raise FileNotFoundError(f"Input Office file not found: {input_path}")

        if not input_path.is_file():
            raise RuntimeError(f"Input path is not a file: {input_path}")

        input_size = input_path.stat().st_size
        if input_size <= 0:
            raise RuntimeError(f"Input Office file is empty: {input_path}")

        suffix = input_path.suffix.lower().strip()

        if suffix not in OFFICE_EXTENSIONS:
            raise RuntimeError(
                f"Unsupported Office file extension: '{suffix}'. "
                f"Supported extensions: {sorted(OFFICE_EXTENSIONS)}. "
                f"Input path: {input_path}"
            )

        libreoffice_path = _resolve_libreoffice_path()

        work_dir_path: Path | None = None

        try:
            work_dir_path = Path(tempfile.mkdtemp(prefix="office_to_pdf_")).resolve()
            input_dir = work_dir_path / "input"
            output_dir = work_dir_path / "output"
            profile_dir = work_dir_path / "lo_profile"

            input_dir.mkdir(parents=True, exist_ok=True)
            output_dir.mkdir(parents=True, exist_ok=True)
            profile_dir.mkdir(parents=True, exist_ok=True)

            safe_input = input_dir / f"input_{uuid.uuid4().hex}{suffix}"
            shutil.copy2(str(input_path), str(safe_input))

            if not safe_input.exists():
                raise RuntimeError(f"Failed to copy input file to temp path: {safe_input}")

            if safe_input.stat().st_size <= 0:
                raise RuntimeError(f"Copied Office file is empty: {safe_input}")

            # LibreOffice on Windows prefers file URI format for profile path.
            profile_uri = profile_dir.as_uri()

            cmd = [
                libreoffice_path,
                "--headless",
                "--nologo",
                "--nofirststartwizard",
                "--nolockcheck",
                "--nodefault",
                "--norestore",
                f"-env:UserInstallation={profile_uri}",
                "--convert-to",
                "pdf",
                "--outdir",
                str(output_dir),
                str(safe_input),
            ]

            logger.info("Office→PDF original input: %s", input_path)
            logger.info("Office→PDF original input size: %s bytes", input_size)
            logger.info("Office→PDF safe input: %s", safe_input)
            logger.info("Office→PDF safe input size: %s bytes", safe_input.stat().st_size)
            logger.info("Office→PDF output dir: %s", output_dir)
            logger.info("Office→PDF LibreOffice path: %s", libreoffice_path)
            logger.info("Office→PDF command: %s", cmd)

            env = os.environ.copy()
            env["HOME"] = str(work_dir_path)
            env["USERPROFILE"] = str(work_dir_path)

            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=180,
                cwd=str(work_dir_path),
                env=env,
            )

            logger.info("LibreOffice return code: %s", result.returncode)
            logger.info("LibreOffice stdout: %s", result.stdout)
            logger.info("LibreOffice stderr: %s", result.stderr)

            if result.returncode != 0:
                raise RuntimeError(
                    "LibreOffice conversion failed\n"
                    f"Return code: {result.returncode}\n"
                    f"Original input: {input_path}\n"
                    f"Original input exists: {input_path.exists()}\n"
                    f"Original input size: {input_path.stat().st_size if input_path.exists() else 'missing'}\n"
                    f"Safe input: {safe_input}\n"
                    f"Safe input exists: {safe_input.exists()}\n"
                    f"Safe input size: {safe_input.stat().st_size if safe_input.exists() else 'missing'}\n"
                    f"Output dir: {output_dir}\n"
                    f"Output files: {[p.name for p in output_dir.glob('*')]}\n"
                    f"STDOUT: {result.stdout}\n"
                    f"STDERR: {result.stderr}"
                )

            expected_pdf = output_dir / f"{safe_input.stem}.pdf"

            if expected_pdf.exists():
                generated_pdf = expected_pdf
            else:
                pdf_files = list(output_dir.glob("*.pdf"))

                if not pdf_files:
                    raise RuntimeError(
                        "LibreOffice finished but did not produce a PDF output\n"
                        f"Expected PDF: {expected_pdf}\n"
                        f"Output dir: {output_dir}\n"
                        f"Output files: {[p.name for p in output_dir.glob('*')]}\n"
                        f"STDOUT: {result.stdout}\n"
                        f"STDERR: {result.stderr}"
                    )

                generated_pdf = pdf_files[0]

            out = _temp_pdf("from_office")
            shutil.copy2(str(generated_pdf), str(out))

            if not out.exists():
                raise RuntimeError(f"Final PDF was not created: {out}")

            if out.stat().st_size <= 0:
                raise RuntimeError(f"Final PDF is empty: {out}")

            logger.info(
                "Office→PDF in %sms → %s, size=%s bytes",
                _ms(t0),
                out,
                out.stat().st_size,
            )

            return out

        except subprocess.TimeoutExpired as exc:
            raise RuntimeError(
                "LibreOffice conversion timed out after 180 seconds"
            ) from exc

        finally:
            if work_dir_path and work_dir_path.exists():
                shutil.rmtree(work_dir_path, ignore_errors=True)

    @staticmethod
    def images_to_pdf(
        image_paths: list[Path],
        page_size: str = "fit",
    ) -> Path:
        """Combine multiple images into a single PDF."""
        t0 = time.time()

        if not image_paths:
            raise RuntimeError("No images were provided for PDF creation")

        doc = fitz.open()

        try:
            for img_path in image_paths:
                img_path = Path(img_path).resolve()

                if not img_path.exists():
                    raise FileNotFoundError(f"Image file not found: {img_path}")

                with Image.open(img_path) as pil_img:
                    w_px, h_px = pil_img.size

                w_pt = w_px * 72 / 96
                h_pt = h_px * 72 / 96

                if page_size == "fit":
                    rect = fitz.Rect(0, 0, w_pt, h_pt)
                else:
                    sizes = {
                        "A4": (595, 842),
                        "Letter": (612, 792),
                        "Legal": (612, 1008),
                        "A3": (842, 1190),
                    }
                    pw, ph = sizes.get(page_size, (595, 842))
                    rect = fitz.Rect(0, 0, pw, ph)

                page = doc.new_page(width=rect.width, height=rect.height)

                scale = min(rect.width / w_pt, rect.height / h_pt)
                img_rect = fitz.Rect(
                    (rect.width - w_pt * scale) / 2,
                    (rect.height - h_pt * scale) / 2,
                    (rect.width + w_pt * scale) / 2,
                    (rect.height + h_pt * scale) / 2,
                )

                page.insert_image(img_rect, filename=str(img_path))

            out = _temp_pdf("from_images")
            doc.save(out, deflate=True)

        finally:
            doc.close()

        logger.info("Images→PDF in %sms, pages=%s → %s", _ms(t0), len(image_paths), out)
        return out

    @staticmethod
    def html_to_pdf(html_content: str) -> Path:
        """Convert HTML string to PDF using PyMuPDF's Story API."""
        t0 = time.time()

        if html_content is None or str(html_content).strip() == "":
            raise RuntimeError("HTML content is empty")

        out = _temp_pdf("from_html")

        story = fitz.Story(html=html_content)
        writer = fitz.DocumentWriter(str(out))

        mediabox = fitz.paper_rect("A4")
        where = mediabox + (36, 36, -36, -36)

        try:
            more = True

            while more:
                device = writer.begin_page(mediabox)
                more, _ = story.place(where)
                story.draw(device)
                writer.end_page()

        finally:
            writer.close()

        if not out.exists():
            raise RuntimeError(f"HTML→PDF failed. Output was not created: {out}")

        logger.info("HTML→PDF in %sms → %s", _ms(t0), out)
        return out